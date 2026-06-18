from pptx import Presentation
import sys

def read_ppt(filepath):
    prs = Presentation(filepath)
    for i, slide in enumerate(prs.slides):
        print(f"=== Slide {i+1} ===")
        for j, shape in enumerate(slide.shapes):
            if hasattr(shape, "text"):
                print(f"Shape {j}: {shape.text}")
        print()

if __name__ == "__main__":
    read_ppt(r"D:\EXONYX\[Pub] ISRO BAH 2026 _ Idea Submission Template.pptx")
