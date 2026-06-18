from pptx import Presentation
from pptx.util import Inches
import os

def update_presentation():
    input_pptx = r"D:\EXONYX\ISRO_BAH_2026_Submission_EXONYX.pptx"
    output_pptx = r"D:\EXONYX\ISRO_BAH_2026_Submission_EXONYX_Final.pptx"
    img_path = r"C:\Users\Aditya Jadhav\.gemini\antigravity\brain\9716afda-a559-4f4e-8d1a-078dbb58bad6\system_architecture_1781715063430.png"
    
    prs = Presentation(input_pptx)
    
    # 1. Update Architecture Slide (Index 6)
    slide7 = prs.slides[6]
    shapes_to_delete = []
    
    for shape in slide7.shapes:
        if hasattr(shape, "text") and "[ Frontend" in shape.text:
            # We will delete this text shape and replace it with the image
            shapes_to_delete.append(shape)
            
    for shape in shapes_to_delete:
        sp = shape._element
        sp.getparent().remove(sp)
        
    # Add the generated image to the slide
    # Position it roughly in the center
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    
    try:
        slide7.shapes.add_picture(img_path, left, top, width=width)
        print("Successfully embedded architecture image.")
    except Exception as e:
        print(f"Error embedding image: {e}")

    # 2. Update Cost Slide (Index 8)
    slide9 = prs.slides[8]
    for shape in slide9.shapes:
        if hasattr(shape, "text") and "Estimated Implementation Cost:" in shape.text:
            shape.text = (
                "Realistic Production Implementation Cost:\n\n"
                "Capital Expenditure (CAPEX):\n"
                "- High-Performance Deep Learning Server (e.g., 1x RTX 6000 Ada or 2x A5000): ~$8,500\n"
                "- High-Speed NAS Storage (50TB for FITS Data Lake): ~$2,500\n"
                "- Total CAPEX: ~$11,000\n\n"
                "Operational Expenditure (OPEX):\n"
                "- Cloud Web Hosting (Dashboard / DB Gateway): ~$1,500 / year\n"
                "- Software Licensing: $0 (Entirely Open Source Stack)\n"
                "- Data Acquisition: $0 (NASA MAST Public Archive)\n"
                "- Total OPEX: ~$1,500 / year\n\n"
                "Conclusion: Highly cost-effective deployment scalable for enterprise/agency-level exoplanet surveying."
            )
            print("Successfully updated production cost.")

    prs.save(output_pptx)
    print(f"Saved final PPT to: {output_pptx}")

if __name__ == "__main__":
    update_presentation()
