from pptx import Presentation
from pptx.util import Inches, Pt
import os

def fill_presentation():
    input_pptx = r"D:\EXONYX\[Pub] ISRO BAH 2026 _ Idea Submission Template.pptx"
    output_pptx = r"D:\EXONYX\ISRO_BAH_2026_Submission_EXONYX.pptx"
    
    prs = Presentation(input_pptx)
    
    # We will iterate through slides and replace the instructional text with our content.
    
    # Slide 3 (Index 2): Opportunity
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if hasattr(shape, "text") and "Opportunity should be able" in shape.text:
            shape.text = (
                "How different is it from existing ideas?\n"
                "Current pipelines rely on expensive cloud computing and massive supercomputer clusters. EXONYX V5 brings research-grade validation to consumer hardware using highly optimized algorithms.\n\n"
                "How will it solve the problem?\n"
                "By combining mathematical Transit Least Squares (TLS) with a custom 1D PyTorch AstroNet CNN, it autonomously isolates and validates transits in noisy Kepler/K2 data without manual intervention.\n\n"
                "USP:\n"
                "A fully containerized, autonomous, and local deep-learning exoplanet discovery pipeline optimized specifically for an RTX 3050 GPU, reducing cloud compute costs to $0."
            )

    # Slide 4 (Index 3): Features
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if hasattr(shape, "text") and "List of features offered" in shape.text:
            shape.text = (
                "Key Features of EXONYX V5:\n"
                "1. Headless Batch Survey Engine: Autonomously crunches thousands of light curves in the background.\n"
                "2. Live Telemetry Dashboard: Next.js UI to monitor metrics, cache sizing, and GPU usage.\n"
                "3. Deep Learning Validation: Integrated PyTorch AstroNet model for False Positive rejection.\n"
                "4. MCMC Characterization: Bayesian 'emcee' integration for precise radius/period uncertainty calculations.\n"
                "5. Automated Reporting: Generates PDF scientific validation reports for every detected candidate."
            )

    # Slide 5 (Index 4): Process Flow
    slide5 = prs.slides[4]
    for shape in slide5.shapes:
        if hasattr(shape, "text") and "Process flow diagram" in shape.text:
            shape.text = (
                "Process Flow (Textual Outline):\n\n"
                "1. Data Ingestion: Download uncalibrated FITS data natively from NASA MAST.\n"
                "2. Detrending: Wōtan filter removes stellar variability and systemic noise.\n"
                "3. TLS Search: Transit Least Squares identifies periodic transit signals.\n"
                "4. Validation: AstroNet1D CNN evaluates the phase-folded curve for False Positive risks.\n"
                "5. Characterization: MCMC walkers sample the posterior distributions for precise parameters.\n"
                "6. Logging: Candidate is saved to the local SQLite DB and broadcast to the UI."
            )

    # Slide 6 (Index 5): Wireframes
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if hasattr(shape, "text") and "Wireframes/Mock diagrams" in shape.text:
            shape.text = (
                "Survey Dashboard UI Components:\n"
                "- Unified Metrics Panel: Displays total candidates, survey progress, and PLI scores.\n"
                "- Live Telemetry: Tracks backend cache sizing and SQLite connection status.\n"
                "- Dark Mode Aesthetics: Premium visual design optimized for data-heavy astronomical workloads.\n"
                "(Note: Actual screenshots can be embedded natively using the 'Insert Image' tool in PowerPoint)."
            )

    # Slide 7 (Index 6): Architecture diagram
    slide7 = prs.slides[6]
    for shape in slide7.shapes:
        if hasattr(shape, "text") and "Architecture diagram" in shape.text:
            shape.text = (
                "System Architecture:\n\n"
                "[ Frontend (Next.js / React) ]\n"
                "          |\n"
                "          v\n"
                "[ API Gateway (FastAPI) ]\n"
                "          |\n"
                "          v\n"
                "[ Core Engine ] -> Wotan Detrender -> TLS -> AstroNet PyTorch CNN -> emcee MCMC\n"
                "          |\n"
                "          v\n"
                "[ Data Layer (SQLite / FITS Cache) ]"
            )

    # Slide 8 (Index 7): Technologies
    slide8 = prs.slides[7]
    for shape in slide8.shapes:
        if hasattr(shape, "text") and "Technologies to be used" in shape.text:
            shape.text = (
                "Technology Stack:\n"
                "- Backend: Python 3.10, FastAPI, Uvicorn.\n"
                "- Astronomy Libraries: Lightkurve, Transit Least Squares, Wōtan, emcee.\n"
                "- Deep Learning: PyTorch (CUDA 11.8 enabled).\n"
                "- Frontend: Next.js, React, Tailwind CSS.\n"
                "- Infrastructure: Docker, Docker Compose, SQLite."
            )

    # Slide 9 (Index 8): Estimated Cost
    slide9 = prs.slides[8]
    for shape in slide9.shapes:
        if hasattr(shape, "text") and "Estimated implementation cost" in shape.text:
            shape.text = (
                "Estimated Implementation Cost:\n\n"
                "- Hardware: $0 (Executes natively on existing local consumer RTX 3050 Laptop GPU).\n"
                "- Software Licensing: $0 (100% open-source stack).\n"
                "- Cloud APIs: $0 (Direct pipeline to NASA MAST public archive).\n"
                "- Maintenance: Negligible (Containerized via Docker for instant reproducibility).\n"
                "Total Cost: $0."
            )

    prs.save(output_pptx)
    print(f"Successfully saved filled presentation to {output_pptx}")

if __name__ == "__main__":
    fill_presentation()
